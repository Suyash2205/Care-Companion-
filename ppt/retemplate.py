"""
retemplate.py — rebuild the deck on the official KJSSE template.

The template carries the real Somaiya crest, the TRUST logo, the three-part footer
bar and the page number on its slide master, at 13.333 x 7.5 in. Our deck was built
at 10 x 5.625 in with a text-box imitation of the letterhead.

Both are 16:9, so every content shape ports across with a single uniform scale of
4/3 — positions, sizes and font sizes alike. Our own letterhead shapes are dropped
because the template's master supplies the real ones.

    python3 retemplate.py
    -> CareCompanion-IE1-KJSSE.pptx
"""
import copy
import glob
import io
import os
import shutil
import zipfile

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))
DECK = os.path.join(HERE, "CareCompanion-IE1.pptx")
TEMPLATE = glob.glob(os.path.expanduser("~/Downloads/KJSSE*PPT*Template*.pptx"))
OUT = os.path.join(HERE, "CareCompanion-IE1-KJSSE.pptx")

SCALE = 4 / 3          # 10 x 5.625 in  ->  13.333 x 7.5 in
FOOTER_TOP = Inches(6.55)    # master footer begins here; content must clear it
CONTENT_TOP = Inches(1.18)   # the crest logo ends at 1.13 in; keep content below it


def is_our_furniture(shape, src_h):
    """Our own letterhead: the header text boxes and the three footer bars.

    Identified by position rather than name, since the shapes are unnamed. The
    template's master replaces all of it.
    """
    top = Emu(shape.top or 0).inches
    return top < 0.72 or top > 4.9


def is_title(shape):
    return (shape.has_text_frame
            and abs((shape.left or 0) - Inches(2.03)) < Inches(0.05)
            and abs((shape.top or 0) - Inches(0.19)) < Inches(0.05))


def scale_shape(el):
    """Scale a shape's geometry and every font size inside it."""
    from pptx.oxml.ns import qn
    # Tables and charts sit in a graphicFrame, whose transform is p:xfrm, not a:xfrm.
    frames = list(el.iter(qn("a:xfrm"))) + list(el.iter(qn("p:xfrm")))
    for xfrm in frames:
        off = xfrm.find(qn("a:off"))
        ext = xfrm.find(qn("a:ext"))
        if off is not None:
            off.set("x", str(int(int(off.get("x")) * SCALE)))
            off.set("y", str(int(int(off.get("y")) * SCALE)))
        if ext is not None:
            ext.set("cx", str(int(int(ext.get("cx")) * SCALE)))
            ext.set("cy", str(int(int(ext.get("cy")) * SCALE)))
    # A table's real geometry lives in its column widths and row heights.
    for col in el.iter(qn("a:gridCol")):
        w = col.get("w")
        if w:
            col.set("w", str(int(int(w) * SCALE)))
    for row in el.iter(qn("a:tr")):
        h = row.get("h")
        if h:
            row.set("h", str(int(int(h) * SCALE)))
    for rpr in el.iter():
        if rpr.tag in (qn("a:rPr"), qn("a:defRPr"), qn("a:endParaRPr")):
            sz = rpr.get("sz")
            if sz:
                rpr.set("sz", str(int(int(sz) * SCALE)))


def clear_logo(el):
    """Nudge a shape below the crest logo.

    The Literature Review and Existing Solutions tables begin at 0.79 in on the old
    canvas, which scales to 1.05 in and runs under the template's logo.
    """
    from pptx.oxml.ns import qn as _qn
    frames = list(el.iter(_qn("a:xfrm"))) + list(el.iter(_qn("p:xfrm")))
    if not frames:
        return
    off = frames[0].find(_qn("a:off"))     # outermost transform only
    if off is not None and int(off.get("y")) < CONTENT_TOP:
        off.set("y", str(int(CONTENT_TOP)))


def port_images(src_slide, new_slide, el):
    """Re-point picture references at the new slide's own relationships.

    Copying a picture's XML alone leaves r:embed pointing at a relationship that
    does not exist in the destination file. python-pptx does not resolve these, so
    the file opens fine there and then fails to render.
    """
    for blip in el.iter(qn("a:blip")):
        rid = blip.get(qn("r:embed"))
        if not rid:
            continue
        blob = src_slide.part.related_part(rid).blob
        # Add by blob, not by relating the part object: reusing the source part keeps
        # its partname (ppt/media/image1.png) and collides with the template's own
        # media of the same name, producing an unopenable package.
        _, new_rid = new_slide.part.get_or_add_image_part(io.BytesIO(blob))
        blip.set(qn("r:embed"), new_rid)


def carry_table_styles(source_deck, dest_deck):
    """Copy the source deck's table styles into the output package.

    The Literature Review and Existing Solutions tables reference a style GUID that
    the template's tableStyles.xml does not define. PowerPoint tolerates the dangling
    reference; LibreOffice refuses to open the file at all.
    """
    with zipfile.ZipFile(source_deck) as z:
        styles = z.read("ppt/tableStyles.xml")
    tmp = dest_deck + ".tmp"
    with zipfile.ZipFile(dest_deck) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = styles if item.filename == "ppt/tableStyles.xml" else zin.read(item.filename)
            zout.writestr(item, data)
    shutil.move(tmp, dest_deck)


def main():
    if not TEMPLATE:
        raise SystemExit("KJSSE template not found in ~/Downloads")
    src = Presentation(DECK)
    out = Presentation(TEMPLATE[0])

    # Strip the template's eight sample slides, keeping its master and layouts.
    lst = out.slides._sldIdLst
    for sid in list(lst):
        rid = sid.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        out.part.drop_rel(rid)
        lst.remove(sid)

    # A layout with no body placeholders, so nothing competes with our shapes.
    layout = next((l for l in out.slide_layouts if l.name == "TITLE_ONLY"),
                  out.slide_layouts[0])

    src_h = src.slide_height
    for s in src.slides:
        new = out.slides.add_slide(layout)
        for sh in list(new.shapes):          # drop the layout's own placeholders
            sh._element.getparent().remove(sh._element)
        for sh in s.shapes:
            if is_our_furniture(sh, src_h) and not is_title(sh):
                continue
            el = copy.deepcopy(sh._element)
            scale_shape(el)
            if not is_title(sh):
                clear_logo(el)
            port_images(s, new, el)
            new.shapes._spTree.append(el)

    out.save(OUT)
    carry_table_styles(DECK, OUT)
    print(f"{len(src.slides)} slides -> {OUT}")
    print(f"canvas {Emu(src.slide_width).inches:.2f}x{Emu(src.slide_height).inches:.2f} "
          f"-> {Emu(out.slide_width).inches:.2f}x{Emu(out.slide_height).inches:.2f} in")


if __name__ == "__main__":
    main()
