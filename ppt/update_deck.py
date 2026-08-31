"""
update_deck.py — adds the slides Internal Evaluation-1 requires, in the deck's own
visual language.

The deck builds slides from hand-positioned shapes on a Somaiya letterhead: white
cards with numbered badges, red header bars over white panels, and a red colour ramp
for flows. New slides reuse those exact geometries and type styles rather than
falling back to plain bullet lists.

    python3 update_deck.py      # needs python-pptx
    -> CareCompanion-IE1.pptx
"""
import copy
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
SRC = os.path.join(HERE, "CareCompanion-source.pptx")
OUT = os.path.join(HERE, "CareCompanion-IE1.pptx")

# ── Brand, lifted from the existing slides ───────────────────────────────────
RED       = RGBColor(0x9B, 0x1C, 0x2E)
RED_DARK  = RGBColor(0x6E, 0x10, 0x1F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
INK       = RGBColor(0x10, 0x10, 0x10)
RAMP      = [RGBColor(0xF1, 0xB7, 0xBE), RGBColor(0xE2, 0x6F, 0x7B),
             RGBColor(0xC4, 0x33, 0x44), RGBColor(0x6E, 0x10, 0x1F)]
HEAD_FONT, BODY_FONT = "Georgia", "Calibri"

TITLE_L, TITLE_T = Inches(2.03), Inches(0.19)
MODEL_TITLE = "Conclusion"


# ── slide plumbing ───────────────────────────────────────────────────────────
def title_box(slide):
    for sh in slide.shapes:
        if sh.has_text_frame and abs(sh.left - TITLE_L) < 6000 and abs(sh.top - TITLE_T) < 6000:
            return sh
    return None


def model_index(prs):
    """Find the model slide by title. A fixed index goes stale as soon as a slide
    is inserted ahead of it, which silently clones the wrong template."""
    for i, s in enumerate(prs.slides):
        tb = title_box(s)
        if tb is not None and tb.text_frame.text.strip() == MODEL_TITLE:
            return i
    raise LookupError(f"no slide titled {MODEL_TITLE!r}")


def is_furniture(sh):
    """Letterhead: the two header marks and the three footer bars."""
    top = Emu(sh.top).inches if sh.top is not None else 0
    return top < 0.75 or top > 4.9


def branded_slide(prs, title):
    """A slide carrying only the letterhead and a title, ready for content."""
    src = prs.slides[model_index(prs)]
    new = prs.slides.add_slide(src.slide_layout)
    for sh in list(new.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        if is_furniture(sh) or sh is title_box(src):
            new.shapes._spTree.append(copy.deepcopy(sh._element))
    tb = title_box(new)
    par = tb.text_frame.paragraphs[0]
    par.runs[0].text = title
    for r in par.runs[1:]:
        r._r.getparent().remove(r._r)
    return new


def move(prs, old_index, new_index):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    lst.remove(ids[old_index])
    lst.insert(new_index, ids[old_index])


# ── drawing helpers, matching the deck's existing geometry ───────────────────
def _text(shape, lines, font, size, color, bold=False, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE, space_after=0):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    for i, line in enumerate(lines):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = align
        par.space_after = Pt(space_after)
        run = par.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def rounded(slide, l, t, w, h, fill):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def card_grid(slide, items, numbered=True, size=12):
    """Two columns of white cards with a dark numbered badge — slides 11/12's pattern."""
    cols, rows = (0.53, 5.25), (1.05, 2.40, 3.75)
    for i, item in enumerate(items[:6]):
        l, t = cols[i % 2], rows[i // 2]
        card = rounded(slide, l, t, 4.20, 1.12, WHITE)
        lines = list(item) if isinstance(item, (list, tuple)) else [item]
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = tf.margin_right = Inches(0.10)
        for j, line in enumerate(lines):
            par = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            par.alignment = PP_ALIGN.CENTER
            run = par.add_run(); run.text = line
            run.font.name = HEAD_FONT if (j == 0 and len(lines) > 1) else BODY_FONT
            run.font.size = Pt(size + 1.5 if (j == 0 and len(lines) > 1) else size)
            run.font.bold = (j == 0 and len(lines) > 1)
            run.font.color.rgb = RED if (j == 0 and len(lines) > 1) else INK
        if numbered:
            badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l + 1.95),
                                           Inches(t - 0.19), Inches(0.38), Inches(0.38))
            badge.fill.solid(); badge.fill.fore_color.rgb = INK
            badge.line.fill.background(); badge.shadow.inherit = False
            _text(badge, [str(i + 1)], HEAD_FONT, 12, WHITE, bold=True)


def stat_cards(slide, items):
    """Cards whose first line is a large figure — for the results slide."""
    cols, rows = (0.53, 5.25), (1.05, 2.40, 3.75)
    for i, (value, label) in enumerate(items[:6]):
        l, t = cols[i % 2], rows[i // 2]
        card = rounded(slide, l, t, 4.20, 1.12, WHITE)
        tf = card.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run(); r1.text = value
        r1.font.name = HEAD_FONT; r1.font.size = Pt(20); r1.font.bold = True
        r1.font.color.rgb = RED
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = label
        r2.font.name = BODY_FONT; r2.font.size = Pt(11); r2.font.color.rgb = INK


def two_panels(slide, left, right, size=12, height=3.15):
    """Red header bar over a white panel, twice — slide 13's pattern."""
    for (l, (heading, lines)) in ((0.53, left), (5.03, right)):
        bar = rounded(slide, l, 1.12, 4.50, 0.53, RED)
        _text(bar, [heading], HEAD_FONT, 16, WHITE, bold=True)
        panel = rounded(slide, l, 1.65, 4.50, height, WHITE)
        panel.text_frame.word_wrap = True
        _text(panel, lines, BODY_FONT, size, INK,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=4)


def chips(slide, labels, top, width=3.00, height=0.60, gap=0.23, size=12, ramp=False):
    """A row of red chips — slide 14's module pattern.

    Solid brand red by default: the pale end of the ramp cannot carry white text.
    """
    n = len(labels)
    total = n * width + (n - 1) * gap
    l = (10.0 - total) / 2
    for i, label in enumerate(labels):
        colour = RAMP[min(i + 1, len(RAMP) - 1)] if ramp else (RED if i % 2 == 0 else RED_DARK)
        chip = rounded(slide, l, top, width, height, colour)
        _text(chip, [label], BODY_FONT, size, WHITE, bold=True)
        l += width + gap


# ── the seven slides ─────────────────────────────────────────────────────────
def slide_technologies(prs):
    s = branded_slide(prs, "Technologies Used")
    card_grid(s, [
        ["Language & UI", "Kotlin · Jetpack Compose · Material 3"],
        ["Architecture", "MVVM · Hilt · Navigation Compose"],
        ["Backend", "Supabase (PostgreSQL) · Row-level security"],
        ["Authentication & Push", "Firebase Auth · Google Sign-In · Cloud Messaging"],
        ["Scheduling", "AlarmManager · WorkManager · Boot receiver"],
        ["Testing & Build", "JUnit · Robolectric · Gradle · Target API 36"],
    ], size=11)
    return s


def slide_status(prs):
    s = branded_slide(prs, "Status of Work Completed")
    two_panels(s,
        ("Features Delivered", [
            "• Google Sign-In replacing phone OTP, for both roles",
            "• Single-use invite codes binding an elder's device",
            "• Guardian: medicines, schedules, reminders, vitals,",
            "   adherence, contacts, family access, SOS monitor",
            "• Elder: home, medicines, contacts, vitals, videos",
            "• Four languages, three text sizes, high contrast",
        ]),
        ("Engineering & Quality", [
            "• Backend live: 17 tables, 33 access policies",
            "• Three-path reminder arming (app, worker, boot)",
            "• Offline adherence queue with later sync",
            "• Fail-honest SOS with user-mediated fallback",
            "• 142 automated tests · Lint clean",
            "• Play-ready: API 36, signed build, no restricted",
            "   permissions",
        ]), size=11)
    return s


def slide_results(prs):
    s = branded_slide(prs, "Implementation Results and Discussion")
    stat_cards(s, [
        ("20 sp vs 13 sp", "Median text size: elder vs guardian interface"),
        ("64 dp minimum", "Every primary control, against a 48 dp standard"),
        ("1 tap", "To raise an emergency alarm"),
        ("3 taps", "To confirm a dose from its reminder"),
        ("142 / 142", "Automated tests passing · Lint clean"),
        ("13 defects, 8 silent", "Found by audit (62%) — all corrected"),
    ])
    return s


def slide_tests(prs):
    s = branded_slide(prs, "Software Test Cases")
    two_panels(s,
        ("Reminders & Adherence", [
            "TC-01  Reminder fires at the scheduled time — Pass",
            "TC-02  Dose taken early cancels its alarm — Pass",
            "TC-03  Alarms restored after restart — Pass",
            "TC-04  Offline dose queued, synced later — Pass",
            "TC-09  Water / walk reminders delivered — Pass",
        ]),
        ("Emergency, Access & UI", [
            "TC-05  SOS reaches guardian; failure stated — Pass",
            "TC-06  Invite code single-use, 7-day expiry — Pass",
            "TC-07  Elder UI correct in all four languages — Pass",
            "TC-08  Save errors shown on network failure — Pass",
            "TC-10  SOS not silenced by notification toggle — Pass",
        ]), size=11, height=2.45)
    chips(s, ["142 automated tests", "JVM + Robolectric", "No emulator required",
              "Lint: 0 errors"], top=4.28, width=2.20, height=0.44, size=10)
    return s


def slide_paper(prs):
    s = branded_slide(prs, "Status of Research Paper")
    two_panels(s,
        ("The Paper", [
            "Title: CareCompanion: Assistive Companion Mobile",
            "Application for Elderly and Disabled Users",
            "",
            "• IEEE format · 7 pages · 6 figures · 21 references",
            "• Contributions: dual-surface design, secure device",
            "   binding, fail-honest alerting, defect taxonomy",
            "• Figures generated from the application source",
        ]),
        ("Submission", [
            "Submitted to ICISS 2027",
            "www.iciss2027.in",
            "",
            "• Peer review managed through Microsoft CMT",
            "• Status: under review",
            "• All references verified against Crossref",
        ]), size=11)
    return s


def slide_future(prs):
    s = branded_slide(prs, "Future Scope")
    card_grid(s, [
        ["Medicine Stock Database", "Track quantity in hand; stock falls as doses are taken"],
        ["Low-Stock Alerts", "Guardian notified before a medicine runs out, to restock in time"],
        ["iOS Application", "Extend CareCompanion to iPhone users"],
        ["Google Play Release", "Publish for public distribution and real-world use"],
        ["User Study", "Evaluate with elderly participants and their caregivers"],
        ["Voice & Vernacular Input", "Speech support for users who find typing difficult"],
    ], size=11)
    return s


def slide_links(prs):
    s = branded_slide(prs, "Supporting Documents and Links")
    card_grid(s, [
        ["Source Code", "github.com/Suyash2205/Care-Companion-"],
        ["Release APK", "github.com/Suyash2205/Care-Companion-/releases"],
        ["Research Paper", "IEEE format — submitted with this presentation"],
        ["User Manual", "PDF — submitted with this presentation"],
        ["Conference Submission", "www.iciss2027.in"],
        ["Demo Video", "To be added before the demonstration"],
    ], size=11)
    return s


# Slide 10 still described the Semester VI prototype.
METHODOLOGY_FIX = {
    "Java": "Kotlin (Jetpack Compose)",
    "Room (Local Storage)": "Supabase / PostgreSQL with row-level security",
    "Android Notifications Manager": "Firebase Cloud Messaging + Notification Manager",
}


def main():
    prs = Presentation(SRC)
    before = len(prs.slides)

    for sh in prs.slides[9].shapes:
        if not sh.has_text_frame:
            continue
        for par in sh.text_frame.paragraphs:
            for run in par.runs:
                if run.text.strip() in METHODOLOGY_FIX:
                    run.text = METHODOLOGY_FIX[run.text.strip()]

    plan = [
        (slide_technologies, 10),
        (slide_status,       24),
        (slide_results,      25),
        (slide_tests,        26),
        (slide_paper,        27),
        (slide_future,       28),
        (slide_links,        37),
    ]
    for builder, target in plan:
        builder(prs)
        move(prs, len(prs.slides) - 1, target)

    prs.save(OUT)
    print(f"{before} slides -> {len(prs.slides)} slides")
    print("Wrote:", OUT)


if __name__ == "__main__":
    main()
